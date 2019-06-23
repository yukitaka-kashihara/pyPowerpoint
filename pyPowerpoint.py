from pptx import Presentation

# コンテンツの定義
name = "柏原 幸隆（かしはら ゆきたか）"
contents = \
"""第２新卒エンジニア採用
趣味：プログラミング、読書
好きな製品：DeepSecurity
好きな技術：Python、Linux
マイブーム：Docker、Djnago、MySQL

こちらのスライドはPythonで作成しました。
ソース：https://github.com/yukitaka-kashihara/pyPowerpoint"""

def main():
    #テンプレートの読み込み
    prs = Presentation('/home/yuki/Downloads/trendTemplate.pptx')

    #テンプレートスライドを取得して、該当する部分にコンテンツを挿入
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text  == 'title':
                        run.text = name
                    
                    if run.text == 'Contents':
                        run.text = contents

    #結果を別ファイルで出力
    prs.save('/home/yuki/Downloads/outputSlide.pptx')

if __name__ == '__main__':
    main()